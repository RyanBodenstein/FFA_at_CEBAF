#!/usr/bin/env python3
"""
Add two thermal aging defense backup slides to the management presentation.
Appends after the current last slide (40).

Slide 41: "Thermal Aging: Quantified and Insufficient"
  - Vendor curve + Arrhenius scaling (left)
  - Excess aging bar chart with 6x shortfall (right)

Slide 42: "Grade Comparison Contradicts Thermal Aging"
  - Dy inversion bar chart (left)
  - Hci sensitivity table (right)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy
import os

PPTX_IN = 'Management_Talk/Bodenstein_Update_June2026_working.pptx'
PPTX_OUT = 'Management_Talk/Bodenstein_Update_June2026_working.pptx'
FIG_DIR = 'Cleanup_Claude/Thermal_Aging_Calc'


def add_title(slide, text, left=0.9, top=0.4, width=11.5, height=0.7, fontsize=32):
    """Add a title text box matching the backup slide style."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fontsize)
    run.font.bold = False
    return txBox


def add_line(slide, left=0.95, top=1.10, width=11.47):
    """Add a horizontal line separator matching backup slide style."""
    from pptx.oxml.ns import qn
    line = slide.shapes.add_connector(
        1,  # straight connector
        Inches(left), Inches(top),
        Inches(left + width), Inches(top))
    # Style the line
    line.line.width = Pt(1)
    return line


def add_slide_number(slide, number, left=10.1, top=6.8, width=3.0, height=0.5):
    """Add slide number in bottom-right corner."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(14)
    return txBox


def add_speaker_notes(slide, text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = text


def main():
    prs = Presentation(PPTX_IN)
    current_count = len(prs.slides)
    print("Current slide count: {}".format(current_count))

    # Use a blank layout (Layout 6: "Body Slide - No Photos 2" has minimal placeholders)
    # But actually let's use a truly blank approach - add from layout and remove placeholders
    # Layout 1 (Transitions Slide) is simplest with just a title placeholder
    blank_layout = prs.slide_layouts[1]  # Transitions Slide

    # =====================================================================
    # SLIDE 41: Thermal Aging - Quantified and Insufficient
    # =====================================================================
    slide41 = prs.slides.add_slide(blank_layout)

    # Remove default placeholders (keep slide clean)
    for ph in list(slide41.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    add_title(slide41, "Thermal Aging: Quantified and Insufficient")
    add_line(slide41)
    add_slide_number(slide41, current_count + 1)

    # Images: two side by side below the line
    # Content area: y starts at ~1.25", available height ~5.5"
    # Each image ~6.0" wide, maintaining aspect ratio
    img_width = 6.0
    img_height = img_width / 1.35  # aspect ratio ~1.35 for both
    img_top = 1.35

    # Left: vendor curve
    slide41.shapes.add_picture(
        os.path.join(FIG_DIR, 'slide_vendor_curve.png'),
        Inches(0.3), Inches(img_top), Inches(img_width), Inches(img_height))

    # Right: excess aging
    slide41.shapes.add_picture(
        os.path.join(FIG_DIR, 'slide_excess_aging.png'),
        Inches(6.8), Inches(img_top), Inches(img_width), Inches(img_height))

    # Key text below images
    text_top = img_top + img_height + 0.15
    txBox = slide41.shapes.add_textbox(
        Inches(0.5), Inches(text_top), Inches(12.3), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True

    points = [
        "Haavisto & Paju (2009): high-Hci NdFeB stable below T\u2080 \u2248 120-150\u00b0C; our magnets at \u226445\u00b0C are far below threshold",
        "Haavisto et al. (2014): 30-year flux loss \u2248 2\u00d7 (1-hour loss) from logarithmic time dependence",
        "Vendor curve is for generic NdFeB (Hci ~12 kOe, max ~80\u00b0C); our grades have 2-4\u00d7 higher coercivity",
    ]
    for i, pt in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = "\u2022 " + pt
        run.font.size = Pt(12)

    # Speaker notes
    add_speaker_notes(slide41,
        "THERMAL AGING DEFENSE - SLIDE 1\n\n"
        "Key talking points:\n"
        "- Dave M. raised the concern that thermal aging (magnetic viscosity/creep) "
        "could explain the -0.208% NdFeB-SmCo differential\n"
        "- We took this seriously and quantified it using Arrhenius scaling from "
        "the vendor curve he provided\n"
        "- LEFT: The vendor curve is for generic NdFeB at 100C. At Ea=1.0 eV, "
        "12 months at 45C equals only ~41 equivalent hours at 100C\n"
        "- RIGHT: The proper comparison is tunnel-vs-lab EXCESS aging (both sets "
        "undergo the same initial knock-down). The excess is 6x too small.\n"
        "- Haavisto T_0 framework (peer-reviewed, IEEE Trans. Magn.): our grades "
        "operate far below T_0 where aging is negligible on decade timescales\n"
        "- Even with the most favorable assumptions (45C for full year, Ea=1.0 eV, "
        "Hci^-2 grade correction), the shortfall is 6x\n\n"
        "References:\n"
        "- Haavisto & Paju (2009) IEEE Trans. Magn. 45(12), 5277-5280\n"
        "- Haavisto et al. (2014) Adv. Mater. Sci. Eng. 2014, 760584\n"
        "- Vendor curve: emagnetsUK.com (generic NdFeB at 100C)"
    )

    print("Added slide {}: Thermal Aging - Quantified and Insufficient".format(current_count + 1))

    # =====================================================================
    # SLIDE 42: Grade Comparison Contradicts Thermal Aging
    # =====================================================================
    slide42 = prs.slides.add_slide(blank_layout)

    for ph in list(slide42.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    add_title(slide42, "Grade Comparison Contradicts Thermal Aging")
    add_line(slide42)
    add_slide_number(slide42, current_count + 2)

    # Left: Dy inversion (aspect ~1.35)
    dy_width = 6.2
    dy_height = dy_width / 1.35
    slide42.shapes.add_picture(
        os.path.join(FIG_DIR, 'slide_dy_inversion.png'),
        Inches(0.2), Inches(1.3), Inches(dy_width), Inches(dy_height))

    # Right: Hci sensitivity table (aspect ~2.3, wider and shorter)
    tbl_width = 6.2
    tbl_height = tbl_width / 2.3
    # Center vertically with the Dy plot
    tbl_top = 1.3 + (dy_height - tbl_height) / 2
    slide42.shapes.add_picture(
        os.path.join(FIG_DIR, 'slide_hci_sensitivity.png'),
        Inches(6.8), Inches(tbl_top), Inches(tbl_width), Inches(tbl_height))

    # Key text below
    text_top2 = 1.3 + dy_height + 0.15
    txBox2 = slide42.shapes.add_textbox(
        Inches(0.5), Inches(text_top2), Inches(12.3), Inches(1.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    points2 = [
        "N42EH (higher Hci, rated 200\u00b0C) degrades MORE than N52SH (lower Hci, rated 150\u00b0C) \u2014 any thermal model predicts the opposite",
        "N42EH contains 1-2% Dy (\u03c3_capture = 994 barns); N52SH contains none \u2014 neutron capture explains the inversion",
        "Even with the most favorable Hci exponent (n=1, linear), the excess aging shortfall is still 2.1\u00d7",
    ]
    for i, pt in enumerate(points2):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = "\u2022 " + pt
        run.font.size = Pt(12)

    # Speaker notes
    add_speaker_notes(slide42,
        "THERMAL AGING DEFENSE - SLIDE 2\n\n"
        "Key talking points:\n"
        "- LEFT: The Dy inversion is the strongest qualitative argument against "
        "pure thermal aging\n"
        "- Thermal aging predicts: lower coercivity = more aging. So N52SH "
        "(Hci=19) should degrade more than N42EH (Hci=30)\n"
        "- We observe the OPPOSITE: N42EH degrades more (-0.252% vs -0.170%)\n"
        "- The difference is 1.6 sigma (suggestive, not definitive on its own)\n"
        "- The natural explanation: N42EH contains 1-2% Dy (994 barn capture "
        "cross-section), concentrated at grain boundaries where damage is most "
        "effective. N52SH has zero Dy.\n"
        "- RIGHT: The Hci correction exponent sensitivity shows the shortfall "
        "persists across all reasonable assumptions\n"
        "- n=1 (linear, most favorable to thermal aging): 2.1x shortfall\n"
        "- n=2 (Stoner-Wohlfarth, our primary estimate): 6.4x shortfall\n"
        "- n=3 (nucleation): 18.5x shortfall\n"
        "- Physical scaling is exponential (Boltzmann factor), giving even larger shortfall\n\n"
        "If pressed: 'We acknowledge the Dy inversion is 1.6 sigma and not "
        "standalone proof. But combined with the 6x quantitative shortfall, "
        "lab controls null, neutron correlation (p=0.03), and degradation in "
        "cool zones, the converging evidence points to radiation, not temperature.'\n\n"
        "References:\n"
        "- Givord et al. (1988) IEEE Trans. Magn. 24(2), 1921-1923 (Hci^2 barrier)\n"
        "- Haavisto et al. (2013) EPJ Web Conf. 40, 06001 (pre-stabilization)"
    )

    print("Added slide {}: Grade Comparison Contradicts Thermal Aging".format(current_count + 2))

    # Save
    prs.save(PPTX_OUT)
    print("\nSaved to: {}".format(PPTX_OUT))
    print("Total slides: {}".format(len(prs.slides)))


if __name__ == '__main__':
    main()
